import os
import collections.abc
import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import openpyxl
import comtypes.client
import sys

base_path = getattr(sys, '_MEIPASS', os.path.dirname(os.path.abspath(__file__)))


class PresentationConverter:
    def __init__(self):
        self.excel_file = None
        self.column1 = None
        self.column2 = None
        self.column3 = None
        self.original_slide = None
        self.output_folder = None
        self.excel_file_label = None
        self.column1_var = None
        self.column1_menu = None
        self.column2_var = None
        self.column2_menu = None
        self.column3_var = None
        self.column3_menu = None
        self.original_slide_label = None
        self.output_folder_label = None
        self.background_image = None

    def add_text_to_slide(self, slide, text, left, top, width, height, font_size=18.0, font_name="Times New Roman",
                          alignment=PP_ALIGN.CENTER):
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        words = text.split(' ')
        p.text = words[0] + '\n' + ' '.join(words[1:])
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.alignment = alignment

    def add_text_to_slide2(self, slide, text, left, top, width, height, font_size=12.8, font_name="Times New Roman",
                           alignment=PP_ALIGN.CENTER):
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        words = text.split(' ')
        p.text = words[0] + ' ' + ' '.join(words[1:])
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.alignment = alignment

    def convert_to_pdf(self):
        if not all([self.excel_file, self.column1, self.column2, self.original_slide, self.output_folder]):
            messagebox.showwarning("Помилка", "Будь ласка, оберіть ексель-файл, стовпці, оригінальну презентацію "
                                              "та папку для збереження.")
            return

        # Завантаження даних з Excel
        workbook = openpyxl.load_workbook(self.excel_file)
        sheet = workbook.active
        data1 = sheet[self.column1]
        data2 = sheet[self.column2]
        data3 = sheet[self.column3]

        # Цикл для обробки даних та створення нових презентацій
        for i, (value1, value2, value3) in enumerate(zip(data1, data2, data3), start=1):
            # Створення нової презентації з оригінальним слайдом
            presentation = Presentation(self.original_slide)

            # Отримання посилання на перший слайд
            slide = presentation.slides[0]
            self.add_text_to_slide(slide, str(value1.value), Pt(380), Pt(70), Pt(40), Pt(20))
            self.add_text_to_slide2(slide, str(value2.value), Pt(80), Pt(340), Pt(40), Pt(20), font_size=16)
            if value3.value == "Жіноча":
                text = "і здобула спеціальність"
            elif value3.value == "Чоловіча":
                text = "і здобув спеціальність"
            else:
                text = " "
            self.add_text_to_slide2(slide, text, Pt(380), Pt(172), Pt(40), Pt(20),)

            output_filename_pptx = os.path.join(self.output_folder, f"presentation_{i}.pptx")

            presentation.save(output_filename_pptx)

        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        for i in range(1, len(data1) + 1):
            output_filename_pptx = os.path.join(self.output_folder, f"presentation_{i}.pptx")
            output_filename_pdf = os.path.join(self.output_folder, f"presentation_{i}.pdf")

            presentation_pdf = powerpoint.Presentations.Open(output_filename_pptx)
            presentation_pdf.ExportAsFixedFormat(output_filename_pdf, 32)  # 32 означає формат PDF
            presentation_pdf.Close()

        powerpoint.Quit()

        messagebox.showinfo("Конвертація завершена", f"Всі презентації успішно сконвертовані у формат PDF. Збережено "
                                                     f"у папці: {self.output_folder}")

    def browse_excel_file(self):
        self.excel_file = filedialog.askopenfilename(filetypes=[("Excel Files", "*.xlsx;*.xls")])
        self.excel_file_label.config(text=os.path.basename(self.excel_file))

    def select_column1(self, *args):
        self.column1 = self.column1_var.get()

    def select_column2(self, *args):
        self.column2 = self.column2_var.get()

    def select_column3(self, *args):
        self.column3 = self.column3_var.get()

    def browse_original_slide(self):
        self.original_slide = filedialog.askopenfilename(filetypes=[("PowerPoint Files", "*.pptx;*.ppt")])
        self.original_slide_label.config(text=os.path.basename(self.original_slide))

    def browse_output_folder(self):
        self.output_folder = filedialog.askdirectory()
        self.output_folder_label.config(text=os.path.basename(self.output_folder))

    def create_gui(self):
        root = tk.Tk()
        current_directory = os.path.dirname(os.path.abspath(__file__))
        # Збільшення розмірів вікна
        root.geometry("800x450")
        root.resizable(width=False, height=False)

        # Фонове зображення
        background_image_path = os.path.join(base_path, "background_image.png")

        if os.path.exists(background_image_path):
            background_image = tk.PhotoImage(file=background_image_path)
            background_label = tk.Label(root, image=background_image)
            background_label.place(x=0, y=0, relwidth=1, relheight=1)

        # Заголовок
        title_label = tk.Label(root, text="Генератор свідоцтв кваліфікованих робітників", font=("Arial", 20, "bold"),
                               highlightthickness=0, bg="#d39684")
        title_label.grid(row=0, column=0, columnspan=4, pady=10)

        # Картинка
        image4_path = os.path.join(base_path, "image4.png")

        if os.path.exists(image4_path):
            image4 = tk.PhotoImage(file=image4_path)
            image_label4 = tk.Label(root, image=image4)
            image_label4.grid(row=1, column=3, rowspan=6, padx=10)

        # Вибір ексель-файлу
        excel_file_button = tk.Button(root, text="Оберіть xl файл", highlightthickness=0,
                                      command=self.browse_excel_file, bg="#b1a491", font=("Arial", 11, "bold"))
        excel_file_button.grid(row=1, column=0, sticky="w", padx=10, pady=10)

        self.excel_file_label = tk.Label(root, text="Файл не обрано", bg="#b1a491")
        self.excel_file_label.grid(row=1, column=1, columnspan=2, sticky="w", padx=10)

        # Вибір першого стовпця
        column1_label = tk.Label(root, text="Стовпець з ПІБ:", highlightthickness=0, bg="#0da6a1",
                                 font=("Arial", 11, "bold"))
        column1_label.grid(row=2, column=0, sticky="w", padx=10, pady=10)

        self.column1_var = tk.StringVar(root)
        self.column1_var.set("F")  # Значення за замовчуванням
        self.column1_var.trace("w", self.select_column1)

        self.column1_menu = tk.OptionMenu(root, self.column1_var, *self.get_excel_columns())
        self.column1_menu.config(bg="#8cb3a1", highlightthickness=0)
        self.column1_menu.grid(row=2, column=1, columnspan=2, sticky="w", padx=10)

        # Вибір другого стовпця
        column2_label = tk.Label(root, text="Стовпець з номером документа:", highlightthickness=0, bg="#0da6a1",
                                 font=("Arial", 11, "bold"))
        column2_label.grid(row=3, column=0, sticky="w", padx=10, pady=10)

        self.column2_var = tk.StringVar(root)
        self.column2_var.set("B")  # Значення за замовчуванням
        self.column2_var.trace("w", self.select_column2)

        self.column2_menu = tk.OptionMenu(root, self.column2_var, *self.get_excel_columns())
        self.column2_menu.config(bg="#8cb3a1", highlightthickness=0)
        self.column2_menu.grid(row=3, column=1, columnspan=2, sticky="w", padx=10)

        # Вибір третього стовпця
        column3_label = tk.Label(root, text="Стовпець де вказана стать", highlightthickness=0, bg="#0da6a1",
                                 font=("Arial", 11, "bold"))
        column3_label.grid(row=4, column=0, sticky="w", padx=10, pady=10)

        self.column3_var = tk.StringVar(root)
        self.column3_var.set("M")  # Значення за замовчуванням
        self.column3_var.trace("w", self.select_column3)

        self.column3_menu = tk.OptionMenu(root, self.column3_var, *self.get_excel_columns())
        self.column3_menu.config(bg="#8cb3a1", highlightthickness=0)
        self.column3_menu.grid(row=4, column=1, columnspan=2, sticky="w", padx=10)

        # Вибір оригінальної презентації
        original_slide_button = tk.Button(root, text="Оберіть шаблон pptx-формату", highlightthickness=0,
                                          command=self.browse_original_slide, bg="#0da6a1", font=("Arial", 11, "bold"))
        original_slide_button.grid(row=5, column=0, sticky="w", padx=10, pady=10)

        self.original_slide_label = tk.Label(root, text="Файл не обрано", bg="#61c2af")
        self.original_slide_label.grid(row=5, column=1, columnspan=2, sticky="w", padx=10)

        # Вибір папки для збереження
        output_folder_button = tk.Button(root, text="Оберіть папку для збереження", highlightthickness=0,
                                         command=self.browse_output_folder, bg="#0da6a1", font=("Arial", 11, "bold"))
        output_folder_button.grid(row=6, column=0, sticky="w", padx=10, pady=10)

        self.output_folder_label = tk.Label(root, text="Папка не обрана", bg="#61c2af")
        self.output_folder_label.grid(row=6, column=1, columnspan=2, sticky="w", padx=10)

        # Кнопки
        convert_button = tk.Button(root, text="Конвертувати в PDF", font=("Arial", 14, "bold"), bg="#007f7c",
                                   fg="white", command=self.convert_to_pdf)
        convert_button.grid(row=8, column=0, pady=20, padx=10)

        exit_button = tk.Button(root, text="Вихід", highlightthickness=0, font=("Arial", 14, "bold"),
                                command=root.quit, padx=10, bg="#1ec8bf")
        exit_button.grid(row=8, column=2, pady=20, padx=20)

        root.mainloop()

    def get_excel_columns(self):
        alphabet = [chr(i) for i in range(ord('A'), ord('Z') + 1)]
        columns = alphabet  # Вибір всіх літер алфавіту від 'A' до 'Z'
        return columns


def main():
    converter = PresentationConverter()
    converter.create_gui()


if __name__ == "__main__":
    main()
